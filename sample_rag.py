# sample_rag.py
import sys
import os
from openai import OpenAI
from dotenv import load_dotenv
import pypdf
from docx import Document
from pptx import Presentation

# 加载 .env 中的 API Key
load_dotenv()
API_KEY = os.getenv("DEEPSEEK_API_KEY")
if not API_KEY:
    print("错误：未找到 DEEPSEEK_API_KEY，请在 .env 文件中设置或设置环境变量")
    sys.exit(1)

client = OpenAI(api_key=API_KEY, base_url="https://api.deepseek.com/v1")

def extract_text(file_path: str) -> str:
    """根据文件扩展名提取文本"""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    elif ext == ".pdf":
        text = ""
        with open(file_path, "rb") as f:
            reader = pypdf.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text
    elif ext == ".docx":
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    elif ext == ".pptx":
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    else:
        raise ValueError(f"不支持的文件类型: {ext}")

def answer_with_rag(doc_content: str, question: str) -> str:
    prompt = f"""请基于以下文档内容回答问题。如果文档中没有相关信息，请说“文档中未提及”。

文档内容：
{doc_content}

问题：{question}
答案："""
    response = client.chat.completions.create(
        model="deepseek-chat",
        messages=[{"role": "user", "content": prompt}],
        temperature=0
    )
    return response.choices[0].message.content

def main():
    # 自动检测文档路径
    if len(sys.argv) >= 2:
        doc_path = sys.argv[1]
        # 如果第二个参数是文件名而非问题，尝试自动检测
        if os.path.exists(doc_path) and len(sys.argv) == 2:
            question = input("请输入问题: ").strip()
        elif len(sys.argv) >= 3:
            question = sys.argv[2]
        else:
            question = input("请输入问题: ").strip()
    else:
        # 默认使用 sample.txt
        if os.path.exists("sample.txt"):
            doc_path = "sample.txt"
            print(f"已自动加载文档: {doc_path}")
            question = input("请输入问题: ").strip()
        else:
            doc_path = input("请输入文档路径: ").strip()
            question = input("请输入问题: ").strip()

    if not os.path.exists(doc_path):
        print(f"错误：文件 {doc_path} 不存在")
        sys.exit(1)

    try:
        print("正在提取文档内容...")
        doc_content = extract_text(doc_path)
        print("正在调用 AI...")
        answer = answer_with_rag(doc_content, question)
        print("\n回答：")
        print(answer)
    except Exception as e:
        print(f"发生错误：{e}")
        sys.exit(1)

if __name__ == "__main__":
    main()