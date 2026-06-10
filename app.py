import streamlit as st
import subprocess
import json
import os
from openai import OpenAI

# 多格式文档解析库
import pypdf
from docx import Document
from pptx import Presentation

# ---------- 配置 ----------
from dotenv import load_dotenv
load_dotenv()

API_KEY = os.getenv("DEEPSEEK_API_KEY")
if not API_KEY:
    st.error("请设置 DEEPSEEK_API_KEY 环境变量或在 .env 文件中配置")
    st.stop()
client = OpenAI(api_key=API_KEY, base_url="https://api.deepseek.com/v1")

# ---------- 多格式文本提取函数 ----------
def extract_text_from_file(uploaded_file) -> str:
    """根据文件扩展名提取文本内容"""
    file_type = uploaded_file.type
    file_name = uploaded_file.name.lower()

    # TXT
    if file_type == "text/plain" or file_name.endswith(".txt"):
        return uploaded_file.read().decode("utf-8", errors="ignore")

    # PDF
    elif file_type == "application/pdf" or file_name.endswith(".pdf"):
        reader = pypdf.PdfReader(uploaded_file)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text

    # Word (DOCX)
    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or file_name.endswith(".docx"):
        doc = Document(uploaded_file)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text

    # PowerPoint (PPTX)
    elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or file_name.endswith(".pptx"):
        prs = Presentation(uploaded_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    else:
        st.error(f"不支持的文件类型: {file_type}")
        return ""

# ---------- 工具函数 ----------
def execute_python(code: str) -> str:
    dangerous_keywords = ["__import__", "exec", "eval", "open", "file", "input", "raw_input"]
    for kw in dangerous_keywords:
        if kw in code:
            return f"错误：代码中包含禁止使用的关键字 '{kw}'"
    try:
        result = subprocess.run(
            ["python", "-c", code],
            capture_output=True,
            text=True,
            timeout=5,
            encoding="utf-8"
        )
        if result.stderr:
            return f"执行错误：{result.stderr}"
        return result.stdout.strip() or "执行成功（无输出）"
    except subprocess.TimeoutExpired:
        return "错误：代码执行超时（>5秒）"
    except Exception as e:
        return f"错误：{str(e)}"

# ---------- Streamlit UI ----------
st.set_page_config(page_title="AI Agent Playground", layout="wide")
st.title("🤖 智能助手（RAG + 代码执行）")
st.markdown("上传 **PDF / Word / PPT / TXT** 文档后提问，我会基于文档回答；如果问题涉及计算，我会自动执行代码。")

if "messages" not in st.session_state:
    st.session_state.messages = []
if "docs_content" not in st.session_state:
    st.session_state.docs_content = ""

with st.sidebar:
    st.header("📄 知识库")
    uploaded_file = st.file_uploader(
        "上传文档（支持 PDF, DOCX, PPTX, TXT）",
        type=["txt", "pdf", "docx", "pptx"]
    )
    if uploaded_file is not None:
        with st.spinner("正在解析文档..."):
            content = extract_text_from_file(uploaded_file)
            if content:
                st.session_state.docs_content = content
                st.success(f"✅ 已加载文档，共 {len(content)} 字符")
                with st.expander("文档预览"):
                    st.text(content[:500] + ("..." if len(content) > 500 else ""))
            else:
                st.error("文档解析失败，请检查格式或内容")

    st.divider()
    if st.button("🗑️ 清空对话"):
        st.session_state.messages = []
        st.session_state.docs_content = ""
        st.rerun()

# 显示历史消息
for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

# 用户输入
if prompt := st.chat_input("问点什么..."):
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    with st.chat_message("assistant"):
        with st.spinner("思考中..."):
            # 构建系统提示：文档上下文 + 工具调用能力
            system_prompt = ""
            if st.session_state.docs_content:
                system_prompt += (
                    "你是一个智能助手。请优先基于以下文档内容回答问题。"
                    "如果问题与文档无关（如数学计算、通用知识），请直接回答或使用工具，"
                    "不要回复\"文档中未提及\"。\n\n"
                    f"文档内容：\n{st.session_state.docs_content}\n\n"
                )
            system_prompt += (
                "如果用户的问题需要执行Python代码（如数学计算、数据处理），"
                "请输出以下JSON格式：\n"
                "{\"tool\": \"execute_python\", \"args\": {\"code\": \"要执行的代码\"}}\n"
                "如果不需要工具，直接输出回答内容。"
            )

            response = client.chat.completions.create(
                model="deepseek-chat",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt}
                ],
                temperature=0
            )
            response_content = response.choices[0].message.content

            # 尝试解析工具调用
            try:
                parsed = json.loads(response_content)
                if parsed.get("tool") == "execute_python":
                    code = parsed["args"]["code"]
                    result = execute_python(code)
                    follow_up = f"代码执行结果：\n{result}\n\n请根据以上结果回答用户问题：{prompt}"
                    final_response = client.chat.completions.create(
                        model="deepseek-chat",
                        messages=[{"role": "user", "content": follow_up}]
                    )
                    answer = final_response.choices[0].message.content
                else:
                    answer = response_content
            except (json.JSONDecodeError, KeyError):
                answer = response_content

        st.markdown(answer)
        st.session_state.messages.append({"role": "assistant", "content": answer})
