"""PPTX 文本提取器，基于 python-pptx 库。"""

from typing import BinaryIO

from pptx import Presentation

from src.backend.engine.parser.base import BaseParser


class PptxParser(BaseParser):
    """使用 python-pptx 提取 PowerPoint 幻灯片中的文本。"""

    def parse(self, file: BinaryIO) -> str:
        prs = Presentation(file)
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return "\n".join(parts)

    @staticmethod
    def supports(extension: str) -> bool:
        return extension == ".pptx"
